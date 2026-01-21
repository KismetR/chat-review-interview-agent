"""
文件处理模块
负责加载和处理各种格式的本地文件
支持两种模式：
1. 普通加载：直接提取文件内容
2. RAG 加载：分块处理，用于向量检索
"""

import os
from pathlib import Path
from typing import Optional, List, Dict, Union


class FileHandler:
    """文件处理器 - 支持 PDF、DOCX、PPTX、TXT、MD"""
    
    SUPPORTED_EXTENSIONS = ['.pdf', '.docx', '.pptx', '.txt', '.md']
    
    def __init__(self, use_unstructured: bool = False, max_content_length: int = 120000):
        """
        Args:
            use_unstructured: 是否使用 unstructured 库（用于 RAG）
            max_content_length: 单个文件最大字符数
        """
        self.use_unstructured = use_unstructured
        self.max_content_length = max_content_length
    
    # ==================== 普通文件加载（不需要 RAG） ====================
    
    def load_files(self, file_paths: Union[str, List[str]]) -> List[Dict]:
        """
        普通文件加载（不分块，直接返回完整内容）
        
        Args:
            file_paths: 文件路径，可以是 str 或 List[str]
            
        Returns:
            list: [{"content": "完整内容", "metadata": {...}}, ...]
            
        Example:
            handler = FileHandler()
            results = handler.load_files("notes.pdf")
            print(results[0]["content"])
        """
        if isinstance(file_paths, str):
            file_paths = [file_paths]
        
        print(f"\n[INFO] 开始加载 {len(file_paths)} 个文件...")
        results = []
        
        for idx, file_path in enumerate(file_paths, 1):
            print(f"[INFO] 正在处理文件 {idx}/{len(file_paths)}: {file_path}")
            
            try:
                result = self._load_single_file(file_path)
                results.append(result)
                print(f"[SUCCESS] ✓ 文件加载成功")
                
            except Exception as e:
                print(f"[ERROR] ✗ 加载失败: {type(e).__name__}: {e}")
                results.append({
                    "content": None,
                    "metadata": {
                        "filename": Path(file_path).name,
                        "error": f"{type(e).__name__}: {e}"
                    }
                })
        
        success_count = sum(1 for r in results if r["content"] is not None)
        print(f"\n[SUMMARY] 加载完成: 成功 {success_count}/{len(file_paths)} 个文件\n")
        
        return results
    
    def _load_single_file(self, file_path: str) -> Dict:
        """加载单个文件的完整内容"""
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        if not path.is_file():
            raise ValueError(f"不是有效的文件: {file_path}")
        
        ext = path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"不支持的文件类型: {ext}")
        
        # 根据类型加载
        content = self._load_by_type(path, ext)
        
        # 限制长度
        original_length = len(content)
        if len(content) > self.max_content_length:
            print(f"[DEBUG] 内容过长 ({original_length} 字符)，截断至 {self.max_content_length}")
            content = content[:self.max_content_length] + "\n\n...(内容过长，已截断)"
        
        return {
            "content": content,
            "metadata": {
                "filename": path.name,
                "type": ext.lstrip('.'),
                "size": path.stat().st_size,
                "char_count": len(content),
                "original_char_count": original_length
            }
        }
    
    def _load_by_type(self, path: Path, ext: str) -> str:
        """根据文件类型加载内容"""
        if ext == '.pdf':
            return self._load_pdf(path)
        elif ext == '.docx':
            return self._load_docx(path)
        elif ext == '.pptx':
            return self._load_pptx(path)
        elif ext in ['.txt', '.md']:
            return self._load_text(path)
        else:
            raise ValueError(f"不支持的类型: {ext}")
    
    def _load_pdf(self, path: Path) -> str:
        """加载 PDF 文件"""
        try:
            import pymupdf
            doc = pymupdf.open(str(path))
            
            if len(doc) == 0:
                raise ValueError("PDF 文件为空")
            
            text_parts = []
            for page_num in range(len(doc)):
                text = doc[page_num].get_text()
                if text.strip():
                    text_parts.append(text)
            
            doc.close()
            
            if not text_parts:
                raise ValueError("PDF 中未找到文本内容")
            
            return "\n\n".join(text_parts)
            
        except ImportError:
            raise ImportError("请安装 PyMuPDF: pip install PyMuPDF")
    
    def _load_docx(self, path: Path) -> str:
        """加载 Word 文件"""
        try:
            from langchain_community.document_loaders import Docx2txtLoader
            loader = Docx2txtLoader(str(path))
            docs = loader.load()
            content = "\n\n".join([doc.page_content for doc in docs])
            
            if not content.strip():
                raise ValueError("DOCX 文件为空")
            
            return content
            
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")
    
    def _load_pptx(self, path: Path) -> str:
        """加载 PowerPoint 文件"""
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_parts.append(f"【幻灯片 {slide_num}】\n" + "\n".join(slide_text))
            
            if not text_parts:
                raise ValueError("PPTX 中未找到文本内容")
            
            return "\n\n".join(text_parts)
            
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")
    
    def _load_text(self, path: Path) -> str:
        """加载纯文本文件"""
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
        
        for encoding in encodings:
            try:
                with open(path, 'r', encoding=encoding) as f:
                    content = f.read()
                
                if not content.strip():
                    raise ValueError("文件为空")
                
                return content
                
            except UnicodeDecodeError:
                continue
        
        raise UnicodeDecodeError(f"无法解码文件 {path.name}")
    
    def format_for_prompt(self, files_data: List[Dict]) -> str:
        """
        将文件数据格式化为 Prompt 文本
        
        Args:
            files_data: load_files() 的返回结果
            
        Returns:
            str: 格式化的文本，可直接插入 Prompt
        """
        if not files_data:
            return "无文档提供。"
        
        formatted_parts = []
        
        for idx, file_data in enumerate(files_data, 1):
            metadata = file_data["metadata"]
            content = file_data["content"]
            
            if content is None:
                formatted_parts.append(
                    f"【文件 {idx}】\n"
                    f"文件名: {metadata['filename']}\n"
                    f"状态: ❌ 加载失败\n"
                    f"错误: {metadata.get('error', '未知错误')}\n"
                )
            else:
                formatted_parts.append(
                    f"【文件 {idx}】\n"
                    f"文件名: {metadata['filename']}\n"
                    f"类型: {metadata['type']}\n"
                    f"大小: {metadata['size']} 字节\n"
                    f"字符数: {metadata['char_count']}\n\n"
                    f"内容:\n{content}\n"
                )
        
        return "\n" + "="*50 + "\n\n".join(formatted_parts) + "\n" + "="*50
    
    # ==================== RAG 文件加载（分块处理） ====================
    
    def load_for_rag(self, file_paths: Union[str, List[str]], chunk_strategy: str = "by_title", max_chars: int = 2000):
        """
        为 RAG 加载文件（分块处理）
        
        Args:
            file_paths: 文件路径
            chunk_strategy: 分块策略 "simple" 或 "by_title"
            max_chars: 每块最大字符数
            
        Returns:
            list: [{"content": "块内容", "metadata": {...}}, ...]
        """
        if self.use_unstructured:
            return self._load_for_rag_unstructured(file_paths, chunk_strategy, max_chars)
        else:
            return self._load_simple_chunks(file_paths, max_chars)
    def _load_simple_chunks(self, file_paths: Union[str, List[str]], max_chars: int) -> List[Dict]:
        """
        简单分块方法（不使用 unstructured）
        按固定字符数分块，有重叠
        """
        if isinstance(file_paths, str):
            file_paths = [file_paths]
        
        all_chunks = []
        chunk_overlap = min(200, max_chars // 10)  # 重叠部分为块大小的10%
        
        # 先用普通方法加载文件
        files_data = self.load_files(file_paths)
        
        for file_data in files_data:
            if file_data["content"] is None:
                print(f"[WARNING] 跳过加载失败的文件: {file_data['metadata']['filename']}")
                continue
            
            content = file_data["content"]
            filename = file_data["metadata"]["filename"]
            
            # 按字符数分块
            chunk_index = 0
            for i in range(0, len(content), max_chars - chunk_overlap):
                chunk_text = content[i:i + max_chars]
                
                if chunk_text.strip():  # 跳过空块
                    all_chunks.append({
                        "content": chunk_text,
                        "metadata": {
                            "source": filename,
                            "chunk_index": chunk_index,
                            "char_start": i,
                            "char_end": min(i + max_chars, len(content))
                        }
                    })
                    chunk_index += 1
        
        print(f"[SUCCESS] 简单分块完成，共生成 {len(all_chunks)} 个块")
        return all_chunks
    
    def _load_for_rag_unstructured(self, file_paths: Union[str, List[str]], chunk_strategy: str, max_chars: int) -> List[Dict]:
        """
        使用 unstructured 库加载文件（高级分块）
        支持按标题、结构等智能分块
        """
        try:
            from unstructured.chunking.title import chunk_by_title
            from unstructured.partition.auto import partition
        except ImportError:
            raise ImportError(
                "请安装 unstructured: pip install unstructured\n"
                "或设置 use_unstructured=False 使用简单分块"
            )
        
        if isinstance(file_paths, str):
            file_paths = [file_paths]
        
        all_chunks = []
        
        for path in file_paths:
            print(f"[DEBUG] 使用 unstructured 处理: {path}")
            
            try:
                # 分区（识别文档结构）
                elements = partition(path, strategy="fast")
                
                # 分块
                if chunk_strategy == "by_title":
                    chunks = chunk_by_title(
                        elements,
                        max_characters=max_chars,
                        combine_text_under_n_chars=100
                    )
                else:
                    chunks = elements
                
                # 转换为标准格式
                for idx, chunk in enumerate(chunks):
                    all_chunks.append({
                        "content": chunk.text,
                        "metadata": {
                            "source": Path(path).name,
                            "chunk_index": idx,
                            "element_type": chunk.category,  # Title, NarrativeText, Table 等
                            "page_number": chunk.metadata.page_number if hasattr(chunk.metadata, "page_number") else None
                        }
                    })
                
                print(f"[SUCCESS] {Path(path).name}: 生成 {len(chunks)} 个结构化块")
                
            except Exception as e:
                print(f"[ERROR] unstructured 处理失败 ({Path(path).name}): {e}")
                print(f"[INFO] 降级使用简单分块方法")
                # 降级到简单分块
                fallback_chunks = self._load_simple_chunks(path, max_chars)
                all_chunks.extend(fallback_chunks)
        
        print(f"[SUMMARY] 总共生成 {len(all_chunks)} 个块")
        return all_chunks


# 便捷函数
def load_and_format_files(file_paths: Union[str, List[str]]) -> str:
    """
    快速加载文件并格式化为 Prompt 文本
    
    Args:
        file_paths: 文件路径
        
    Returns:
        str: 格式化的文件内容
    """
    handler = FileHandler()
    files_data = handler.load_files(file_paths)
    return handler.format_for_prompt(files_data)